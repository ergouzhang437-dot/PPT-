"""
Create a demo PPT with AI-generated images using python-pptx.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

IMG_DIR = os.path.join(os.path.dirname(__file__), "ppt_images")
OUT = os.path.join(os.path.dirname(__file__), "AI_PPT_Demo.pptx")

# Color palette (Ocean theme)
PRIMARY = RGBColor(0x06, 0x5A, 0x82)
SECONDARY = RGBColor(0x1C, 0x72, 0x93)
DARK = RGBColor(0x21, 0x29, 0x5C)
LIGHT = RGBColor(0xCA, 0xDC, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x64, 0x74, 0x8B)
BG_LIGHT = RGBColor(0xF8, 0xF9, 0xFA)

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)


def add_full_rect(slide, color, x=0, y=0, w=10, h=5.625):
    """Add a full-slide background rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, font_size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, font_face="Calibri"):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_face
    p.alignment = align
    return txBox


def add_image(slide, path, x, y, w, h, **kwargs):
    """Add an image."""
    return slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))


# ══════════════════════════════════════════════════════
# Slide 1: Title
# ═══════════════════════════════════════════════════════
s1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_full_rect(s1, DARK)
add_image(s1, os.path.join(IMG_DIR, "ai_28fd34c0_1280x1024.png"), 0, 0, 10, 5.625)
# Semi-transparent overlay
ov = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(5.625))
ov.fill.solid()
ov.fill.fore_color.rgb = DARK
ov.fill.fore_color.brightness = 0.3
ov.line.fill.background()

add_text(s1, "AI Image Generation", 1, 1.2, 8, 1, 44, True, WHITE, PP_ALIGN.LEFT, "Arial Black")
add_text(s1, "Automated Visual Content for Presentations", 1, 2.4, 8, 0.6, 18, False, LIGHT, PP_ALIGN.LEFT, "Calibri")


# ═══════════════════════════════════════════════════════
# Slide 2: What We Built
# ═══════════════════════════════════════════════════════
s2 = prs.slides.add_slide(prs.slide_layouts[6])
add_full_rect(s2, BG_LIGHT)
# Accent bar
bar = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.12), Inches(5.625))
bar.fill.solid()
bar.fill.fore_color.rgb = PRIMARY
bar.line.fill.background()

add_text(s2, "What We Built", 0.5, 0.4, 5, 0.6, 32, True, PRIMARY, PP_ALIGN.LEFT, "Arial Black")

# Bullet text
txBox = s2.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5), Inches(2))
tf = txBox.text_frame
tf.word_wrap = True

items = [
    "Image Generation Module",
    "Python script calling DashScope WAN 2.7 API",
    "Supports Chinese & English prompts",
    "Outputs PNG + base64 for PPT embedding",
]
for i, item in enumerate(items):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(15)
    p.font.bold = (i == 0)
    p.font.name = "Calibri"
    p.font.color.rgb = MUTED
    p.space_after = Pt(6)

# Image on right
add_image(s2, os.path.join(IMG_DIR, "ai_1b383e49_1280x1024.png"), 5.5, 0.6, 4, 3.2)


# ═══════════════════════════════════════════════════════
# Slide 3: The Pipeline
# ═══════════════════════════════════════════════════════
s3 = prs.slides.add_slide(prs.slide_layouts[6])
add_full_rect(s3, PRIMARY)

add_text(s3, "The Pipeline", 1, 0.3, 8, 0.6, 36, True, WHITE, PP_ALIGN.LEFT, "Arial Black")

steps = [
    ("01", "Generate", "AI generates image\nfrom text prompt", RGBColor(0x1C, 0x72, 0x93)),
    ("02", "Process", "Saves to local file\n+ base64 encoding", RGBColor(0x21, 0x29, 0x5C)),
    ("03", "Embed", "Insert into slide\nvia pptxgenjs", RGBColor(0x02, 0xC3, 0x9A)),
]
card_w, card_h = 2.8, 3.5
start_x = 0.4
for i, (num, title, desc, color) in enumerate(steps):
    x = start_x + i * (card_w + 0.3)

    # Card
    card = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.2), Inches(card_w), Inches(card_h))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.fill.background()

    # Number circle
    circ = s3.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.25), Inches(1.45), Inches(0.6), Inches(0.6))
    circ.fill.solid()
    circ.fill.fore_color.rgb = color
    circ.line.fill.background()
    add_text(s3, num, x + 0.25, 1.45, 0.6, 0.6, 20, True, WHITE, PP_ALIGN.CENTER, "Arial Black")

    add_text(s3, title, x + 0.25, 2.2, card_w - 0.5, 0.5, 22, True, PRIMARY, PP_ALIGN.LEFT, "Arial Black")
    add_text(s3, desc, x + 0.25, 2.8, card_w - 0.5, 1.5, 13, False, MUTED, PP_ALIGN.LEFT, "Calibri")


# ═══════════════════════════════════════════════════════
# Slide 4: Example Result
# ═══════════════════════════════════════════════════════
s4 = prs.slides.add_slide(prs.slide_layouts[6])
add_full_rect(s4, BG_LIGHT)
bar2 = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.12), Inches(5.625))
bar2.fill.solid()
bar2.fill.fore_color.rgb = PRIMARY
bar2.line.fill.background()

add_text(s4, "Example: AI-Generated Image", 0.5, 0.3, 5, 0.6, 28, True, PRIMARY, PP_ALIGN.LEFT, "Arial Black")
add_text(s4, 'Prompt: "一间有着精致窗户的花店，\n漂亮的木质门，摆放着花朵"',
         0.5, 1.1, 4.5, 1.2, 14, False, MUTED, PP_ALIGN.LEFT, "Calibri")

add_image(s4, os.path.join(IMG_DIR, "ai_fa5c8fa6_1280x1024.png"), 0.5, 2.5, 4.2, 2.8)
add_image(s4, os.path.join(IMG_DIR, "ai_1ffa6977_1280x1024.png"), 5.3, 1.0, 4.2, 4.2)
add_text(s4, "AI Generated", 5.3, 5.3, 4.2, 0.3, 11, False, MUTED, PP_ALIGN.CENTER, "Calibri")


# ═══════════════════════════════════════════════════════
# Slide 5: Closing
# ═══════════════════════════════════════════════════════
s5 = prs.slides.add_slide(prs.slide_layouts[6])
add_full_rect(s5, DARK)
add_image(s5, os.path.join(IMG_DIR, "ai_1ffa6977_1280x1024.png"), 0, 0, 10, 5.625)
ov2 = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(5.625))
ov2.fill.solid()
ov2.fill.fore_color.rgb = DARK
ov2.fill.fore_color.brightness = 0.3
ov2.line.fill.background()

add_text(s5, "AI + PPT", 1, 1.5, 8, 1.2, 52, True, WHITE, PP_ALIGN.CENTER, "Arial Black")
add_text(s5, "Every slide can have unique AI-generated visuals", 1, 2.9, 8, 0.5, 16, False, LIGHT, PP_ALIGN.CENTER, "Calibri")


prs.save(OUT)
print(f"PPT saved: {OUT}")
